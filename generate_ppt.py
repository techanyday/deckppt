import os
import logging
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_LINE, MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from apis.openai_api import OpenAIClient
import re
import tempfile
from io import BytesIO
from datetime import datetime

class ColorPalette:
    """Modern color palettes for professional presentations."""
    MINIMALIST_BLUE = {
        "background": RGBColor(255, 255, 255),  # White
        "shape": RGBColor(230, 240, 250),      # Light Blue
        "text": RGBColor(44, 62, 80),          # Dark Navy
        "title": RGBColor(41, 128, 185),        # Blue
        "primary": RGBColor(41, 128, 185),        # Blue
        "secondary": RGBColor(44, 62, 80),          # Dark Navy
        "accent": RGBColor(230, 240, 250)      # Light Blue
    }
    
    SOFT_GRAY = {
        "background": RGBColor(255, 255, 255),  # White
        "shape": RGBColor(242, 242, 242),      # Light Gray
        "text": RGBColor(51, 51, 51),          # Dark Gray
        "title": RGBColor(44, 62, 80),          # Dark Blue-Gray
        "primary": RGBColor(44, 62, 80),          # Dark Blue-Gray
        "secondary": RGBColor(51, 51, 51),          # Dark Gray
        "accent": RGBColor(242, 242, 242)      # Light Gray
    }
    
    FRESH_GREEN = {
        "background": RGBColor(255, 255, 255),  # White
        "shape": RGBColor(230, 247, 230),      # Light Green
        "text": RGBColor(46, 139, 87),         # Dark Green
        "title": RGBColor(39, 174, 96),         # Green
        "primary": RGBColor(39, 174, 96),         # Green
        "secondary": RGBColor(46, 139, 87),         # Dark Green
        "accent": RGBColor(230, 247, 230)      # Light Green
    }
    
    ELEGANT_PURPLE = {
        "background": RGBColor(250, 250, 250),  # Light Gray
        "shape": RGBColor(239, 230, 250),      # Soft Lavender
        "text": RGBColor(75, 0, 130),          # Dark Purple
        "title": RGBColor(142, 68, 173),        # Purple
        "primary": RGBColor(142, 68, 173),        # Purple
        "secondary": RGBColor(75, 0, 130),          # Dark Purple
        "accent": RGBColor(239, 230, 250)      # Soft Lavender
    }
    
    PROFESSIONAL_TEAL = {
        "background": RGBColor(255, 255, 255),  # White
        "shape": RGBColor(230, 250, 247),      # Light Teal
        "text": RGBColor(0, 128, 128),         # Dark Teal
        "title": RGBColor(22, 160, 133),        # Teal
        "primary": RGBColor(22, 160, 133),        # Teal
        "secondary": RGBColor(0, 128, 128),         # Dark Teal
        "accent": RGBColor(230, 250, 247)      # Light Teal
    }
    
    @classmethod
    def get_palette(cls, theme="minimalist_blue"):
        """Get color palette by theme name."""
        palettes = {
            "minimalist_blue": cls.MINIMALIST_BLUE,
            "soft_gray": cls.SOFT_GRAY,
            "fresh_green": cls.FRESH_GREEN,
            "elegant_purple": cls.ELEGANT_PURPLE,
            "professional_teal": cls.PROFESSIONAL_TEAL
        }
        return palettes.get(theme, cls.MINIMALIST_BLUE)

def create_shaped_textbox(slide, left, top, width, height, text, palette, 
                         is_title=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Create a shaped textbox with modern styling."""
    # Add shape background
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = palette["shape"]
    shape.line.fill.background()  # No outline
    
    # Add text
    text_box = slide.shapes.add_textbox(
        left + Inches(0.25),  # Add padding
        top + Inches(0.25),
        width - Inches(0.5),
        height - Inches(0.5)
    )
    
    frame = text_box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    p = frame.add_paragraph()
    p.text = text
    p.font.size = Pt(24 if is_title else 18)
    p.font.name = 'Calibri'
    p.font.bold = is_title
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = palette["text"]
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    
    return shape, text_box

def create_modern_content_slide(ppt, title, insights, palette):
    """Create a modern content slide with shaped text blocks."""
    layout = ppt.slide_layouts[6]  # Blank layout
    slide = ppt.slides.add_slide(layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]
    
    # Add title shape with gradient fill
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(12.33), Inches(0.8)
    )
    
    # Apply gradient fill to title shape
    fill = title_shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = palette["primary"]
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[1].color.rgb = palette["secondary"]
    fill.gradient_stops[1].position = 1
    title_shape.line.width = 0
    
    # Add title text
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.6),
        Inches(11.83), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.margin_bottom = 0
    title_frame.margin_left = 0
    
    p = title_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.name = 'Segoe UI Light'
    p.font.bold = False
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Calculate layout for insights
    content_top = Inches(1.7)
    content_bottom = Inches(7.0)
    available_height = content_bottom - content_top
    
    if len(insights) <= 3:
        # Single row layout
        block_width = Inches(3.8)
        block_height = Inches(4.8)
        total_width = block_width * len(insights)
        spacing = (Inches(12.33) - total_width) / (len(insights) + 1)
        
        for i, insight in enumerate(insights):
            left = Inches(0.5) + spacing + (block_width + spacing) * i
            
            # Add shape background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, content_top,
                block_width, block_height
            )
            
            # Shape styling
            shape.fill.solid()
            shape.fill.fore_color.rgb = palette["shape"]
            shape.line.color.rgb = palette["accent"]
            shape.line.width = Pt(1)
            
            # Add text
            text_box = slide.shapes.add_textbox(
                left + Inches(0.3),
                content_top + Inches(0.3),
                block_width - Inches(0.6),
                block_height - Inches(0.6)
            )
            
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = tf.add_paragraph()
            p.text = insight
            p.font.size = Pt(16)
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.LEFT
            p.font.color.rgb = palette["text"]
            p.space_before = Pt(0)
            p.space_after = Pt(12)
            p.line_spacing = 1.2
    else:
        # 2x2 grid layout
        block_width = Inches(5.67)
        block_height = Inches(2.4)
        h_spacing = Inches(0.5)
        v_spacing = Inches(0.4)
        
        for i, insight in enumerate(insights[:4]):
            row = i // 2
            col = i % 2
            
            left = Inches(0.5) + (block_width + h_spacing) * col
            top = content_top + (block_height + v_spacing) * row
            
            # Add shape background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top,
                block_width, block_height
            )
            
            # Shape styling
            shape.fill.solid()
            shape.fill.fore_color.rgb = palette["shape"]
            shape.line.color.rgb = palette["accent"]
            shape.line.width = Pt(1)
            
            # Add text
            text_box = slide.shapes.add_textbox(
                left + Inches(0.3),
                top + Inches(0.3),
                block_width - Inches(0.6),
                block_height - Inches(0.6)
            )
            
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = tf.add_paragraph()
            p.text = insight
            p.font.size = Pt(16)
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.LEFT
            p.font.color.rgb = palette["text"]
            p.space_before = Pt(0)
            p.space_after = Pt(6)
            p.line_spacing = 1.2
    
    return slide

def create_modern_conclusion_slide(ppt, key_insights, palette):
    """Create a modern conclusion slide with shaped takeaways."""
    layout = ppt.slide_layouts[6]  # Blank layout
    slide = ppt.slides.add_slide(layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(11.33), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.add_paragraph()
    p.text = "Key Takeaways"
    p.font.size = Pt(40)
    p.font.name = 'Calibri Light'
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = palette["title"]
    
    # Add insights in shaped boxes with equal spacing
    block_width = Inches(11.33)
    block_height = Inches(1.2)
    v_spacing = Inches(0.3)
    
    for i, insight in enumerate(key_insights[:4]):  # Limit to 4 takeaways
        top = Inches(1.8) + (block_height + v_spacing) * i
        
        # Add shape background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), top, block_width, block_height
        )
        
        # Shape styling
        shape.fill.solid()
        shape.fill.fore_color.rgb = palette["shape"]
        shape.line.color.rgb = palette["shape"]
        shape.line.width = Pt(1)
        
        # Add text
        text_box = slide.shapes.add_textbox(
            Inches(1.25), top + Inches(0.2),
            block_width - Inches(0.5), block_height - Inches(0.4)
        )
        
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        
        p = tf.add_paragraph()
        p.text = insight
        p.font.size = Pt(18)
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = palette["text"]
        p.space_before = Pt(0)
        p.space_after = Pt(6)
    
    return slide

def create_title_slide(ppt, title, palette):
    """Create a clean, modern title slide."""
    layout = ppt.slide_layouts[6]  # Blank layout
    slide = ppt.slides.add_slide(layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]
    
    # Add title with modern styling
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), 
        Inches(10.33), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(54)
    p.font.name = 'Calibri'
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = palette["title"]
    
    return slide

def generate_intro_slide(ppt, title, palette):
    """Generate a modern introduction slide."""
    layout = ppt.slide_layouts[6]  # Blank layout
    slide = ppt.slides.add_slide(layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), 
        Inches(11.33), Inches(0.8)
    )
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "Overview"
    p.font.size = Pt(36)
    p.font.name = 'Calibri'
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = palette["title"]
    
    # Generate and add overview text in a shaped box
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key:
        raise ValueError("OPENAI_API_KEY environment variable is not set")
        
    client = OpenAIClient(api_key)
    prompt = f"""Write a compelling 2-3 sentence introduction for {title}.
    Requirements:
    1. Start with a powerful market insight or trend
    2. Focus on business impact and opportunity
    3. Keep it concise (40-60 words)
    4. Use active voice and present tense
    5. NO phrases like 'this presentation' or 'we will discuss'
    
    Example:
    'The renewable energy sector is experiencing unprecedented growth, with global investments exceeding $500B in 2024. Advanced technologies and favorable policies are accelerating adoption, creating new opportunities for businesses to lead in sustainability while reducing operational costs.'"""
    
    try:
        overview_text = client.generate(prompt).strip()
    except Exception as e:
        logging.error(f"Error generating overview: {e}")
        overview_text = f"The {title.lower()} landscape is rapidly evolving, presenting unprecedented opportunities for innovation and growth. Organizations that embrace these changes and implement strategic solutions will gain significant competitive advantages in the coming years."
    
    create_shaped_textbox(
        slide, Inches(1), Inches(1.8),
        Inches(11.33), Inches(2),
        overview_text, palette
    )
    
    return slide

def get_theme_layout_ids(theme="minimalist_blue"):
    """Get the layout and color IDs for the selected theme."""
    return {
        "colors": ColorPalette.get_palette(theme)
    }

def apply_theme_color(paragraph, color):
    """Apply theme color to paragraph text."""
    paragraph.font.color.rgb = color

def create_presentation(topic, num_slides=5, theme="minimalist_blue"):
    """Create a modern, professional presentation."""
    ppt = Presentation()
    palette = ColorPalette.get_palette(theme)
    
    # Set slide size to widescreen
    ppt.slide_width = Inches(13.33)
    ppt.slide_height = Inches(7.5)
    
    # Title slide
    create_title_slide(ppt, topic, palette)
    
    # Overview slide
    generate_intro_slide(ppt, topic, palette)
    
    # Generate insights
    insights = generate_content_sections(topic, (num_slides - 2) * 3)  # -2 for title and overview
    
    # Create content slides (one slide per 3 insights)
    for i in range(0, len(insights), 3):
        slide_insights = insights[i:i+3]
        if slide_insights:
            create_modern_content_slide(
                ppt,
                f"Key Insights",
                slide_insights,
                palette
            )
    
    # Ensure we have the exact number of slides requested
    current_slides = len(ppt.slides)
    if current_slides < num_slides:
        # Generate additional insights if needed
        additional_insights = generate_content_sections(topic, (num_slides - current_slides) * 3)
        
        # Create remaining slides
        for i in range(0, len(additional_insights), 3):
            if len(ppt.slides) < num_slides:
                slide_insights = additional_insights[i:i+3]
                if slide_insights:
                    create_modern_content_slide(
                        ppt,
                        f"Additional Insights",
                        slide_insights,
                        palette
                    )
    
    return ppt

def generate_content_sections(topic, num_sections):
    """Generate unique content sections without numbering."""
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key:
        raise ValueError("OPENAI_API_KEY environment variable is not set")
        
    client = OpenAIClient(api_key)
    
    prompt = f"""Create {num_sections} distinct insights about {topic} for a modern business presentation.
    Each insight should be a complete thought that fits in a small text block (30-40 words).
    
    Requirements:
    1. Focus on business impact and strategic value
    2. Include specific metrics or examples
    3. Start with action verbs
    4. Be forward-looking and actionable
    5. No bullet points or lists
    6. Each insight must be unique (no repetition)
    
    Example insights:
    "Implement AI-powered customer analytics to increase retention by 25% through personalized engagement strategies and predictive behavior modeling."
    
    "Deploy blockchain-based supply chain tracking to reduce operational costs by 40% while ensuring end-to-end transparency and compliance."
    
    Format: Return each insight as a separate paragraph."""
    
    try:
        response = client.generate(prompt)
        insights = [insight.strip() for insight in response.strip().split("\n\n")
                   if insight.strip()][:num_sections]
        
        return insights
    except Exception as e:
        logging.error(f"Error generating insights: {e}")
        return []

def generate_ppt(topic, num_slides=5, theme="minimalist_blue"):
    """Generate a professional presentation."""
    # Clean the topic for file naming
    clean_topic = re.sub(r'[^\w\s-]', '', topic.replace('/', '_'))
    
    try:
        # Create presentation with modern design
        ppt = create_presentation(topic, num_slides, theme)
        
        # Save to a temporary file with a secure name
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        temp_filename = f"{clean_topic}_{timestamp}.pptx"
        temp_path = os.path.join(tempfile.gettempdir(), temp_filename)
        
        ppt.save(temp_path)
        logging.info(f"Presentation saved to {temp_path}")
        
        return temp_path
        
    except Exception as e:
        logging.error(f"Error generating presentation: {e}")
        raise
